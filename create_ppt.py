from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()

# 设置幻灯片母版标题样式
def set_title_format(shape):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(139, 69, 19)

def set_content_format(shape):
    tf = shape.text_frame
    for p in tf.paragraphs:
        p.font.size = Pt(24)
        p.line_spacing = 1.5

# 幻灯片 1: 封面
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "母爱如山"
subtitle.text = "论母爱的伟大与无私\n\n关于母亲的文章"

set_title_format(title)
set_content_format(subtitle)

# 幻灯片 2: 引言
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "引言"
content.text = "世界上的一切光荣和骄傲，都来自母亲。-- 高尔基\n\n" \
              "母亲，这个平凡而伟大的称谓，承载着人类最深沉的情感。\n\n" \
              "母爱如涓涓细流，滋润着我们的心田；\n" \
              "母爱如巍峨高山，为我们遮风挡雨。"

set_title_format(title)
set_content_format(content)

# 幻灯片 3: 母爱的无私与奉献
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "一、母爱的无私与奉献"
content.text = "谁言寸草心，报得三春晖。-- 孟郊《游子吟》\n\n" \
              "母亲对子女的爱：\n" \
              "- 不求回报\n" \
              "- 不计得失\n" \
              "- 无私奉献\n\n" \
              "母亲啊！你是荷叶，我是红莲。-- 冰心"

set_title_format(title)
set_content_format(content)

# 幻灯片 4: 母爱的细腻与包容
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "二、母爱的细腻与包容"
content.text = "母爱的特征：\n\n" \
              "- 细腻：敏锐察觉子女情绪变化\n" \
              "- 关怀：生病时的悉心照料\n" \
              "- 安慰：挫折时的温柔鼓励\n\n" \
              "母爱是一种巨大的火焰。-- 罗曼罗兰\n\n" \
              "这火焰不仅温暖，更能照亮子女前行的道路。"

set_title_format(title)
set_content_format(content)

# 幻灯片 5: 母爱对人格塑造的影响
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "三、母爱对人格塑造的影响"
content.text = "心理学观点：\n\n" \
              "弗洛伊德认为：\n" \
              "早期母婴关系对个体心理发展具有决定性作用\n\n" \
              "得到充分母爱的孩子：\n" \
              "- 更加自信\n" \
              "- 更加乐观\n" \
              "- 具有更强的安全感\n\n" \
              "教育不能创造什么，但它能启发儿童创造力 -- 陶行知"

set_title_format(title)
set_content_format(content)

# 幻灯片 6: 现代社会中的母爱
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "四、现代社会中的母爱"
content.text = "现代母爱的多元化表现：\n\n" \
              "- 职业母亲：平衡工作与家庭\n" \
              "- 单亲母亲：独自承担养育责任\n\n" \
              "现代母爱应有之义：\n" \
              "- 适当的放手与引导\n" \
              "- 让子女学会独立\n\n" \
              "所谓父女母子一场，只不过意味着，你和他的缘分就是\n" \
              "今生今世不断地在目送他的背影渐行渐远。-- 龙应台《目送》"

set_title_format(title)
set_content_format(content)

# 幻灯片 7: 结语
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "结语"
content.text = "慈母手中线，游子身上衣。\n\n" \
              "母亲的爱，如同那细密的针脚，\n" \
              "编织着我们生命的底色。\n\n" \
              "母亲的爱是永恒的，它超越时间、空间和生死的界限。-- 乔治麦克唐纳\n\n" \
              "让我们用实际行动回报母爱，\n" \
              "让这份人间最珍贵的情感得以传承与延续。"

set_title_format(title)
set_content_format(content)

# 幻灯片 8: 参考文献
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "参考文献"
content.text = "[1] 高尔基。母亲 [M].北京：人民文学出版社，2005.\n\n" \
              "[2] 冰心。冰心散文选 [M].北京：人民文学出版社，1997.\n\n" \
              "[3] 孟郊。游子吟 [A].全唐诗 [C].北京：中华书局，1960.\n\n" \
              "[4] 罗曼罗兰。约翰克利斯朵夫 [M].北京：人民文学出版社，2000.\n\n" \
              "[5] 陶行知。陶行知教育文集 [M].成都：四川教育出版社，2005.\n\n" \
              "[6] 龙应台。目送 [M].北京：生活读书新知三联书店，2009.\n\n" \
              "[7] 弗洛伊德。精神分析引论 [M].北京：商务印书馆，1984.\n\n" \
              "[8] 乔治麦克唐纳。北风的背后 [M].上海：上海译文出版社，2007."

set_title_format(title)
tf = content.text_frame
for p in tf.paragraphs:
    p.font.size = Pt(16)
    p.line_spacing = 1.3

# 保存 PPT
prs.save('C:/Users/28054/.openclaw/workspace/母爱如山 - 关于母亲的文章.pptx')
print('PPT 已创建成功！')
print('文件位置：C:\\Users\\28054\\.openclaw\\workspace\\母爱如山 - 关于母亲的文章.pptx')
print('')
print('幻灯片结构：')
print('1. 封面 - 母爱如山')
print('2. 引言 - 高尔基名言')
print('3. 母爱的无私与奉献')
print('4. 母爱的细腻与包容')
print('5. 母爱对人格塑造的影响')
print('6. 现代社会中的母爱')
print('7. 结语')
print('8. 参考文献')
